from fastapi import UploadFile
from pdf2image import convert_from_path
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
import pandas as pd
import os
from pathlib import Path

from starlette.responses import FileResponse

OUTPUT_DIR = "static/converted_files"
os.makedirs(OUTPUT_DIR, exist_ok=True)

async def convert_pdf_to_jpg(file: UploadFile):
    input_path = save_upload_file(file)
    images = convert_from_path(input_path)
    output_files = []
    for i, image in enumerate(images):
        output_path = os.path.join(OUTPUT_DIR, f"{Path(file.filename).stem}_{i + 1}.jpg")
        image.save(output_path, "JPEG")
        output_files.append(output_path)
    return {"files": output_files}

async def convert_pdf_to_word(file: UploadFile):
    input_path = save_upload_file(file)
    reader = PdfReader(input_path)
    output_path = input_path.replace(".pdf", ".docx")
    doc = Document()
    for page in reader.pages:
        doc.add_paragraph(page.extract_text())
    doc.save(output_path)
    return FileResponse(output_path)

async def convert_pdf_to_excel(file: UploadFile):
    input_path = save_upload_file(file)
    reader = PdfReader(input_path)
    output_path = input_path.replace(".pdf", ".xlsx")
    data = []
    for page in reader.pages:
        data.append([line for line in page.extract_text().splitlines()])
    df = pd.DataFrame(data)
    df.to_excel(output_path, index=False, header=False)
    return FileResponse(output_path)

async def convert_pdf_to_ppt(file: UploadFile):
    input_path = save_upload_file(file)
    reader = PdfReader(input_path)
    output_path = input_path.replace(".pdf", ".pptx")
    ppt = Presentation()
    for page in reader.pages:
        slide = ppt.slides.add_slide(ppt.slide_layouts[5])
        textbox = slide.shapes.add_textbox(left=0, top=0, width=720, height=540)
        textbox.text = page.extract_text()
    ppt.save(output_path)
    return FileResponse(output_path)

def save_upload_file(file: UploadFile):
    input_path = os.path.join(OUTPUT_DIR, file.filename)
    with open(input_path, "wb") as f:
        f.write(file.file.read())
    return input_path

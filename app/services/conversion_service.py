import base64
import tempfile
import os
import uuid
from io import BytesIO
from fpdf import FPDF
from PIL import Image
import img2pdf
from pdf2image import convert_from_path
from docx import Document
from pptx import Presentation
import pandas as pd
import pdfplumber
from PyPDF2 import PdfReader

# Define a folder to store converted files
SAVE_FOLDER = "converted_files"
os.makedirs(SAVE_FOLDER, exist_ok=True)  # Ensure folder exists

def encode_io_to_base64(file_io: BytesIO) -> str:
    """Encodes a BytesIO object to a Base64 string."""
    file_io.seek(0)
    return base64.b64encode(file_io.read()).decode("utf-8")


async def jpg_to_pdf(input_io: BytesIO) -> BytesIO:
    """Converts JPG or PNG to PDF."""
    try:
        output_io = BytesIO()
        output_io.write(img2pdf.convert(input_io))
        output_io.seek(0)
        return output_io
    except Exception as e:
        raise ValueError(f"JPG-to-PDF conversion failed: {e}")


async def word_to_pdf(input_io: BytesIO) -> BytesIO:
    try:
        input_io.seek(0)

        # Create a temporary file
        with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as temp_file:
            temp_file.write(input_io.read())
            temp_file_path = temp_file.name

        # Load document from the temporary file path
        doc = Document(temp_file_path)

        # Generate PDF
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)

        for paragraph in doc.paragraphs:
            text = paragraph.text.replace('–', '-').replace('—', '-')
            text = text.encode('latin-1', 'ignore').decode('latin-1')
            pdf.multi_cell(0, 10, txt=text)

        output_io = BytesIO()
        pdf.output(output_io)
        output_io.seek(0)

        # Cleanup the temporary file
        os.remove(temp_file_path)

        return output_io
    except Exception as e:
        if 'temp_file_path' in locals():
            os.remove(temp_file_path)  # Ensure cleanup on error
        raise ValueError(f"Word-to-PDF conversion failed: {e}")

class PDF(FPDF):
    def header(self):
        """Repeat table headers on each new page."""
        if hasattr(self, 'table_headers') and hasattr(self, 'column_widths'):
            self.set_font("Arial", style="B", size=9)
            for col_name, col_width in zip(self.table_headers, self.column_widths):
                self.cell(col_width, 10, col_name, border=1, align="C")
            self.ln()

    def footer(self):
        """Add footer with page number."""
        self.set_y(-15)
        self.set_font("Arial", size=8)
        self.cell(0, 10, f'Page {self.page_no()}', align="C")

async def excel_to_pdf(input_io: BytesIO) -> str:
    """Converts an Excel spreadsheet to a dynamically sized PDF."""
    try:
        input_io.seek(0)  # Reset file pointer

        # Read Excel file
        df = pd.read_excel(input_io, engine="openpyxl")

        # Generate a unique filename
        pdf_filename = os.path.join(SAVE_FOLDER, "converted_file.pdf")

        # Estimate required PDF width based on column count
        base_column_width = 30  # Minimum column width in mm
        calculated_width = len(df.columns) * base_column_width
        min_width = 210  # Minimum width (A4)
        max_width = 420  # Maximum width (A3)
        page_width = min(max(calculated_width, min_width), max_width)  # Scale within limits

        # Create a PDF with a dynamically determined width
        pdf = PDF(orientation="L", unit="mm", format=(page_width, 297))  # Adjust width dynamically
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()
        pdf.set_font("Arial", size=9)

        # Calculate column widths dynamically
        column_widths = [max(pdf.get_string_width(str(col)) + 10, base_column_width) for col in df.columns]

        # Store headers for the header function
        pdf.table_headers = df.columns.tolist()
        pdf.column_widths = column_widths

        # Write table headers
        pdf.set_font("Arial", style="B", size=9)
        for col_name, col_width in zip(df.columns, column_widths):
            pdf.cell(col_width, 10, col_name, border=1, align="C")
        pdf.ln()

        # Write table rows
        pdf.set_font("Arial", size=8)
        row_height = 8  # Standard row height

        for _, row in df.iterrows():
            row_max_height = row_height  # Track max height per row
            y_before = pdf.get_y()

            for col, col_width in zip(row, column_widths):
                x_before = pdf.get_x()
                pdf.multi_cell(col_width, row_height, str(col), border=1, align="C")  # Wrap text
                row_max_height = max(row_max_height, pdf.get_y() - y_before)
                pdf.set_xy(x_before + col_width, y_before)

            pdf.ln(row_max_height)

            # Ensure new page is added if needed
            if pdf.get_y() > 190:
                pdf.add_page()

        pdf.output(pdf_filename)  # Save the file

        return pdf_filename  # Return saved file path

    except Exception as e:
        raise ValueError(f"Excel-to-PDF conversion failed: {e}")

async def ppt_to_pdf(input_io: BytesIO) -> BytesIO:
    """Converts PowerPoint presentations to PDF."""
    try:
        output_io = BytesIO()
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        ppt = Presentation(input_io)
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    pdf.multi_cell(0, 10, txt=shape.text)
        pdf.output(output_io)
        output_io.seek(0)
        return output_io
    except Exception as e:
        raise ValueError(f"PPT-to-PDF conversion failed: {e}")


async def html_to_pdf(input_io: BytesIO) -> BytesIO:
    """Converts HTML content to PDF."""
    try:
        output_io = BytesIO()
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        for line in input_io.read().decode("utf-8", errors="replace").splitlines():
            pdf.multi_cell(0, 10, txt=line.strip())
        pdf.output(output_io)
        output_io.seek(0)
        return output_io
    except Exception as e:
        raise ValueError(f"HTML-to-PDF conversion failed: {e}")


async def pdf_to_jpg(input_io: BytesIO):
    """Converts PDF pages to JPG images."""
    try:
        with tempfile.NamedTemporaryFile(delete=True, suffix=".pdf") as temp_pdf:
            temp_pdf.write(input_io.read())
            temp_pdf.flush()

            images = convert_from_path(temp_pdf.name)
            output_files = []
            for image in images:
                output_io = BytesIO()
                image.save(output_io, "JPEG")
                output_io.seek(0)
                output_files.append(encode_io_to_base64(output_io))

        return output_files
    except Exception as e:
        raise ValueError(f"PDF-to-JPG conversion failed: {e}")


async def pdf_to_word(input_io: BytesIO) -> BytesIO:
    """Converts PDF to editable Word document."""
    try:
        output_io = BytesIO()
        doc = Document()

        with pdfplumber.open(input_io) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    doc.add_paragraph(text)

        doc.save(output_io)
        output_io.seek(0)
        return output_io
    except Exception as e:
        raise ValueError(f"PDF-to-Word conversion failed: {e}")


async def pdf_to_ppt(input_io: BytesIO) -> BytesIO:
    """Converts PDF pages into PowerPoint slides."""
    try:
        output_io = BytesIO()
        reader = PdfReader(input_io)
        ppt = Presentation()
        for page in reader.pages:
            text = page.extract_text()
            if text:
                slide = ppt.slides.add_slide(ppt.slide_layouts[5])
                textbox = slide.shapes.add_textbox(left=0, top=0, width=720, height=540)
                textbox.text = text
        ppt.save(output_io)
        output_io.seek(0)
        return output_io
    except Exception as e:
        raise ValueError(f"PDF-to-PPT conversion failed: {e}")


async def pdf_to_excel(input_io: BytesIO) -> BytesIO:
    """Extracts table data from PDF and converts it to an Excel spreadsheet."""
    try:
        output_io = BytesIO()
        data = []

        with pdfplumber.open(input_io) as pdf:
            for page in pdf.pages:
                tables = page.extract_tables()
                for table in tables:
                    for row in table:
                        data.append(row)

        df = pd.DataFrame(data)
        df.to_excel(output_io, index=False, header=False)
        output_io.seek(0)
        return output_io
    except Exception as e:
        raise ValueError(f"PDF-to-Excel conversion failed: {e}")

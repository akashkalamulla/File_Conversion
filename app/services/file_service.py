from fastapi import UploadFile, HTTPException
from fpdf import FPDF
from PIL import Image
import os
import img2pdf
from docx import Document
from pptx import Presentation
import pandas as pd
from pathlib import Path

from starlette.responses import FileResponse

UPLOAD_DIR = "static/uploaded_files"
OUTPUT_DIR = "static/converted_files"
os.makedirs(UPLOAD_DIR, exist_ok=True)
os.makedirs(OUTPUT_DIR, exist_ok=True)

async def convert_file_to_pdf(file_type: str, file: UploadFile):
    if file_type == "jpg" and file.content_type.startswith("image/"):
        return await jpg_to_pdf(file)
    if file_type == "word" and file.filename.endswith(".docx"):
        return await word_to_pdf(file)
    if file_type == "excel" and file.filename.endswith(".xlsx"):
        return await excel_to_pdf(file)
    if file_type == "ppt" and file.filename.endswith(".pptx"):
        return await ppt_to_pdf(file)
    raise HTTPException(status_code=400, detail=f"Unsupported file type: {file_type}")

async def jpg_to_pdf(file: UploadFile):
    input_path = save_upload_file(file)
    output_path = input_path.replace(Path(input_path).suffix, ".pdf")
    with open(output_path, "wb") as pdf_file:
        pdf_file.write(img2pdf.convert(input_path))
    return FileResponse(output_path)

async def word_to_pdf(file: UploadFile):
    input_path = save_upload_file(file)
    output_path = input_path.replace(".docx", ".pdf")
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    doc = Document(input_path)
    for paragraph in doc.paragraphs:
        pdf.cell(200, 10, txt=paragraph.text, ln=True)
    pdf.output(output_path)
    return FileResponse(output_path)

async def excel_to_pdf(file: UploadFile):
    input_path = save_upload_file(file)
    output_path = input_path.replace(".xlsx", ".pdf")
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    df = pd.read_excel(input_path)
    for _, row in df.iterrows():
        pdf.cell(200, 10, txt=str(row.values), ln=True)
    pdf.output(output_path)
    return FileResponse(output_path)

async def ppt_to_pdf(file: UploadFile):
    input_path = save_upload_file(file)
    output_path = input_path.replace(".pptx", ".pdf")
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    ppt = Presentation(input_path)
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                pdf.cell(200, 10, txt=shape.text, ln=True)
    pdf.output(output_path)
    return FileResponse(output_path)

def save_upload_file(file: UploadFile):
    input_path = os.path.join(UPLOAD_DIR, file.filename)
    with open(input_path, "wb") as f:
        f.write(file.file.read())
    return input_path
